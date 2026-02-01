#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "Pillow"]
# ///
"""
Create presentations using existing PPTX templates.
Preserves master slides, layouts, and branding from the template.
"""

import argparse
import json
import copy
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SKILL_DIR = Path(__file__).parent.parent
TEMPLATES_DIR = SKILL_DIR / "templates"


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def get_layout_by_name(prs: Presentation, name: str):
    """Find a layout by name (case-insensitive partial match)."""
    name_lower = name.lower()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if name_lower in layout.name.lower():
                return layout
    # Default to first content layout
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


def create_from_template(template_path: str, slides_data: list, output_path: str, 
                          clear_content: bool = True, start_from: int = 0):
    """
    Create a new presentation based on an existing template.
    
    Args:
        template_path: Path to the template PPTX
        slides_data: List of slide definitions
        output_path: Where to save the result
        clear_content: If True, remove existing slides except master
        start_from: Keep slides 0 to start_from-1 from template
    """
    prs = Presentation(template_path)
    
    # Optionally clear existing content slides but keep structure
    if clear_content:
        # Keep only the first N slides (title, navigation, etc.)
        slides_to_keep = start_from
        while len(prs.slides) > slides_to_keep:
            rId = prs.slides._sldIdLst[slides_to_keep].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[slides_to_keep]
    
    # Add new slides from data
    for slide_def in slides_data:
        layout_name = slide_def.get("layout", "Title and Content")
        layout = get_layout_by_name(prs, layout_name)
        
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if slide.shapes.title and slide_def.get("title"):
            slide.shapes.title.text = slide_def["title"]
        
        # Set content/bullets
        if slide_def.get("bullets"):
            # Find content placeholder
            for shape in slide.placeholders:
                if "content" in shape.name.lower() and shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    
                    for i, bullet in enumerate(slide_def["bullets"]):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = bullet
                        p.level = slide_def.get("bullet_level", 0)
                    break
        
        # Add custom text box
        if slide_def.get("text_box"):
            tb = slide_def["text_box"]
            txBox = slide.shapes.add_textbox(
                Inches(tb.get("left", 1)),
                Inches(tb.get("top", 1)),
                Inches(tb.get("width", 5)),
                Inches(tb.get("height", 1))
            )
            tf = txBox.text_frame
            tf.text = tb.get("text", "")
            tf.word_wrap = True
        
        # Add rounded rectangle (navigation button style)
        if slide_def.get("nav_buttons"):
            for btn in slide_def["nav_buttons"]:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(btn.get("left", 1)),
                    Inches(btn.get("top", 1)),
                    Inches(btn.get("width", 2)),
                    Inches(btn.get("height", 0.8))
                )
                shape.text = btn.get("text", "")
                
                if btn.get("fill_color"):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb(btn["fill_color"])
                
                if btn.get("line_color"):
                    shape.line.color.rgb = hex_to_rgb(btn["line_color"])
        
        # Add image from path
        if slide_def.get("image"):
            img = slide_def["image"]
            img_path = img.get("path")
            if img_path and Path(img_path).exists():
                slide.shapes.add_picture(
                    img_path,
                    Inches(img.get("left", 7)),
                    Inches(img.get("top", 1.5)),
                    width=Inches(img.get("width", 5))
                )
        
        # Speaker notes
        if slide_def.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_def["notes"]
    
    prs.save(output_path)
    return len(prs.slides)


def main():
    parser = argparse.ArgumentParser(description="Create presentation from template")
    parser.add_argument("--template", "-t", required=True, help="Template PPTX path or name")
    parser.add_argument("--slides", "-s", required=True, help="JSON file with slide definitions")
    parser.add_argument("--output", "-o", default="output.pptx", help="Output path")
    parser.add_argument("--keep-slides", "-k", type=int, default=0, 
                       help="Keep first N slides from template")
    parser.add_argument("--append", "-a", action="store_true",
                       help="Append to existing slides instead of clearing")
    
    args = parser.parse_args()
    
    # Find template
    template_path = args.template
    if not Path(template_path).exists():
        # Check templates directory
        template_path = TEMPLATES_DIR / f"{args.template}.pptx"
        if not template_path.exists():
            print(f"Template not found: {args.template}")
            print(f"Available templates in {TEMPLATES_DIR}:")
            for t in TEMPLATES_DIR.glob("*.pptx"):
                print(f"  - {t.stem}")
            return
    
    # Load slide definitions
    with open(args.slides, "r") as f:
        slides_data = json.load(f)
    
    if not isinstance(slides_data, list):
        slides_data = slides_data.get("slides", [])
    
    # Create presentation
    num_slides = create_from_template(
        str(template_path),
        slides_data,
        args.output,
        clear_content=not args.append,
        start_from=args.keep_slides
    )
    
    print(f"Created: {args.output}")
    print(f"Total slides: {num_slides}")


if __name__ == "__main__":
    main()
