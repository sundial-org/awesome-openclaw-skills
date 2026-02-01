#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx"]
# ///
"""
Analyze a PowerPoint template to extract layouts, colors, fonts, and structure.
Use this to understand existing templates before creating new presentations.
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def analyze_template(pptx_path: str, verbose: bool = False) -> dict:
    """Analyze a PowerPoint template and return its structure."""
    prs = Presentation(pptx_path)
    
    analysis = {
        "file": str(pptx_path),
        "dimensions": {
            "width_inches": round(prs.slide_width.inches, 2),
            "height_inches": round(prs.slide_height.inches, 2),
        },
        "slide_count": len(prs.slides),
        "masters": [],
        "layouts": [],
        "slides": [],
    }
    
    # Analyze slide masters
    for i, master in enumerate(prs.slide_masters):
        master_info = {
            "index": i,
            "layout_count": len(master.slide_layouts),
            "layouts": []
        }
        
        for j, layout in enumerate(master.slide_layouts):
            layout_info = {
                "index": j,
                "name": layout.name,
                "placeholders": []
            }
            
            for ph in layout.placeholders:
                layout_info["placeholders"].append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
            
            master_info["layouts"].append(layout_info)
            analysis["layouts"].append({
                "name": layout.name,
                "index": j,
                "master": i,
            })
        
        analysis["masters"].append(master_info)
    
    # Analyze slides (sample or all)
    max_slides = len(prs.slides) if verbose else min(10, len(prs.slides))
    
    for idx in range(max_slides):
        slide = prs.slides[idx]
        slide_info = {
            "number": idx + 1,
            "layout": slide.slide_layout.name,
            "shapes": []
        }
        
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type) if shape.shape_type else "Unknown",
            }
            
            # Get position/size
            shape_info["left"] = round(shape.left.inches, 2) if shape.left else 0
            shape_info["top"] = round(shape.top.inches, 2) if shape.top else 0
            shape_info["width"] = round(shape.width.inches, 2) if shape.width else 0
            shape_info["height"] = round(shape.height.inches, 2) if shape.height else 0
            
            # Get text content
            if hasattr(shape, 'text') and shape.text.strip():
                shape_info["text"] = shape.text.strip()[:100]
            
            # Get fill color
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type is not None:
                        shape_info["fill_type"] = str(shape.fill.type)
                        if hasattr(shape.fill, 'fore_color'):
                            try:
                                rgb = shape.fill.fore_color.rgb
                                if rgb:
                                    shape_info["fill_color"] = f"#{rgb}"
                            except:
                                pass
                except:
                    pass
            
            # Get font info from first paragraph
            if hasattr(shape, 'text_frame'):
                try:
                    for para in shape.text_frame.paragraphs[:1]:
                        if para.runs:
                            run = para.runs[0]
                            font = run.font
                            font_info = {}
                            if font.name:
                                font_info["name"] = font.name
                            if font.size:
                                font_info["size"] = font.size.pt
                            if font.bold:
                                font_info["bold"] = True
                            if font.color.rgb:
                                font_info["color"] = f"#{font.color.rgb}"
                            if font_info:
                                shape_info["font"] = font_info
                except:
                    pass
            
            slide_info["shapes"].append(shape_info)
        
        analysis["slides"].append(slide_info)
    
    return analysis


def print_summary(analysis: dict):
    """Print a human-readable summary."""
    print(f"📊 Template Analysis: {Path(analysis['file']).name}")
    print(f"   Dimensions: {analysis['dimensions']['width_inches']}\" x {analysis['dimensions']['height_inches']}\"")
    print(f"   Total slides: {analysis['slide_count']}")
    print()
    
    print("📐 Available Layouts:")
    for layout in analysis["layouts"]:
        print(f"   [{layout['index']}] {layout['name']}")
    print()
    
    print("📄 Slide Structure (first 10):")
    for slide in analysis["slides"][:10]:
        print(f"\n   Slide {slide['number']}: {slide['layout']}")
        for shape in slide["shapes"][:5]:
            text_preview = f" → \"{shape.get('text', '')[:40]}\"" if shape.get('text') else ""
            print(f"      • {shape['type']}: {shape['name']}{text_preview}")
        if len(slide["shapes"]) > 5:
            print(f"      ... and {len(slide['shapes']) - 5} more shapes")


def main():
    parser = argparse.ArgumentParser(description="Analyze PowerPoint template")
    parser.add_argument("template", help="Path to PPTX file")
    parser.add_argument("--json", "-j", action="store_true", help="Output JSON")
    parser.add_argument("--verbose", "-v", action="store_true", help="Analyze all slides")
    parser.add_argument("--output", "-o", help="Save JSON to file")
    
    args = parser.parse_args()
    
    analysis = analyze_template(args.template, verbose=args.verbose)
    
    if args.json or args.output:
        output = json.dumps(analysis, indent=2)
        if args.output:
            with open(args.output, "w") as f:
                f.write(output)
            print(f"Saved analysis to {args.output}")
        else:
            print(output)
    else:
        print_summary(analysis)


if __name__ == "__main__":
    main()
