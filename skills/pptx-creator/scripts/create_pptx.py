#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "httpx", "pyyaml", "Pillow"]
# ///
"""
PowerPoint Creator - Generate professional presentations from outlines, topics, or data.
"""

import argparse
import json
import os
import re
import sys
import tempfile
from pathlib import Path
from datetime import datetime
from io import BytesIO

import httpx
import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)

# Skill directory (for templates)
SKILL_DIR = Path(__file__).parent.parent
TEMPLATES_DIR = SKILL_DIR / "templates"

# Style presets - these match create_template.py PRESETS
STYLES = {
    "minimal": {
        "title_font": "Helvetica Neue",
        "body_font": "Helvetica Neue",
        "title_size": 44,
        "body_size": 18,
        "title_color": "1a1a1a",
        "body_color": "4a4a4a",
        "accent_color": "0066cc",
    },
    "corporate": {
        "title_font": "Arial",
        "body_font": "Arial",
        "title_size": 40,
        "body_size": 18,
        "title_color": "003366",
        "body_color": "333333",
        "accent_color": "0066cc",
    },
    "creative": {
        "title_font": "Avenir Next",
        "body_font": "Avenir",
        "title_size": 48,
        "body_size": 20,
        "title_color": "2d2d2d",
        "body_color": "4a4a4a",
        "accent_color": "ff5722",
    },
    "dark": {
        "title_font": "SF Pro Display",
        "body_font": "SF Pro Text",
        "title_size": 44,
        "body_size": 18,
        "title_color": "FFFFFF",
        "body_color": "e0e0e0",
        "accent_color": "00d4ff",
    },
    "executive": {
        "title_font": "Georgia",
        "body_font": "Calibri",
        "title_size": 42,
        "body_size": 18,
        "title_color": "1e3a5f",
        "body_color": "2c3e50",
        "accent_color": "c9a227",
    },
    "startup": {
        "title_font": "Poppins",
        "body_font": "Inter",
        "title_size": 48,
        "body_size": 20,
        "title_color": "2d3436",
        "body_color": "636e72",
        "accent_color": "6c5ce7",
    },
}


def parse_outline(outline_path: str) -> dict:
    """Parse markdown outline into slide structure."""
    with open(outline_path, "r") as f:
        content = f.read()
    
    presentation = {
        "title": "",
        "subtitle": "",
        "author": "",
        "slides": []
    }
    
    lines = content.strip().split("\n")
    current_slide = None
    
    for line in lines:
        line = line.rstrip()
        
        # Presentation title (# Title)
        if line.startswith("# ") and not presentation["title"]:
            presentation["title"] = line[2:].strip()
            continue
        
        # Metadata (subtitle:, author:)
        if line.lower().startswith("subtitle:"):
            presentation["subtitle"] = line.split(":", 1)[1].strip()
            continue
        if line.lower().startswith("author:"):
            presentation["author"] = line.split(":", 1)[1].strip()
            continue
        
        # Slide header (## Slide N: Title or ## Title)
        if line.startswith("## "):
            if current_slide:
                presentation["slides"].append(current_slide)
            
            title = line[3:].strip()
            # Remove "Slide N:" prefix if present
            title = re.sub(r"^Slide\s+\d+:\s*", "", title, flags=re.IGNORECASE)
            
            current_slide = {
                "title": title,
                "layout": "title_and_content",
                "bullets": [],
                "notes": "",
                "image": None,
                "chart": None,
                "table": None,
                "data_source": None,
            }
            continue
        
        # Slide content
        if current_slide:
            # Bullet point
            if line.startswith("- "):
                item = line[2:].strip()
                
                # Check for special directives
                if item.lower().startswith("chart:"):
                    current_slide["chart"] = item.split(":", 1)[1].strip()
                    current_slide["layout"] = "chart"
                elif item.lower().startswith("table:"):
                    current_slide["table"] = item.split(":", 1)[1].strip()
                    current_slide["layout"] = "table"
                elif item.lower().startswith("data:") or item.lower().startswith("source:"):
                    current_slide["data_source"] = item.split(":", 1)[1].strip()
                elif item.lower().startswith("layout:"):
                    current_slide["layout"] = item.split(":", 1)[1].strip()
                elif item.startswith("!["):
                    # Image: ![alt](path or generate: prompt)
                    match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", item)
                    if match:
                        alt, src = match.groups()
                        current_slide["image"] = {"alt": alt, "src": src.strip()}
                        if "generate:" not in src.lower():
                            current_slide["layout"] = "image_and_text"
                else:
                    current_slide["bullets"].append(item)
            
            # Speaker notes (indented or after ---)
            elif line.strip().startswith(">"):
                current_slide["notes"] += line.strip()[1:].strip() + "\n"
    
    # Don't forget the last slide
    if current_slide:
        presentation["slides"].append(current_slide)
    
    return presentation


def generate_image(prompt: str, output_path: str) -> str | None:
    """Generate image using nano-banana-pro skill."""
    try:
        import subprocess
        script_path = SKILL_DIR.parent / "nano-banana-pro" / "scripts" / "generate_image.py"
        
        if not script_path.exists():
            print(f"Warning: nano-banana-pro not found, skipping image generation", file=sys.stderr)
            return None
        
        result = subprocess.run(
            ["uv", "run", str(script_path), "--prompt", prompt, "--filename", output_path, "--resolution", "1K"],
            capture_output=True,
            text=True,
            timeout=60
        )
        
        if result.returncode == 0:
            # Parse output for file path
            for line in result.stdout.split("\n"):
                if line.startswith("Image saved:"):
                    return line.split(":", 1)[1].strip()
            return output_path
        else:
            print(f"Warning: Image generation failed: {result.stderr}", file=sys.stderr)
            return None
    except Exception as e:
        print(f"Warning: Image generation error: {e}", file=sys.stderr)
        return None


def fetch_crm_data(query: str) -> list[dict]:
    """Fetch data from CRM."""
    api_url = os.environ.get("TWENTY_API_URL", "")
    api_token = os.environ.get("TWENTY_API_TOKEN", "")
    
    if not api_url or not api_token:
        print("Warning: CRM not configured, skipping data fetch", file=sys.stderr)
        return []
    
    # Parse query like "twenty://opportunities?stage=negotiation"
    # For now, return empty - this would need proper CRM API integration
    return []


def create_presentation(data: dict, style: str = "minimal", template_path: str | None = None) -> Presentation:
    """Create PowerPoint presentation from structured data."""
    
    # Start with template or blank
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
    
    style_config = STYLES.get(style, STYLES["minimal"])
    
    # Title slide
    if data.get("title"):
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        if title:
            title.text = data["title"]
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(style_config["title_size"])
                paragraph.font.name = style_config["title_font"]
                paragraph.font.bold = True
                paragraph.font.color.rgb = hex_to_rgb(style_config["title_color"])
        
        if data.get("subtitle") and len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = data["subtitle"]
            if data.get("author"):
                subtitle.text += f"\n{data['author']}"
            for paragraph in subtitle.text_frame.paragraphs:
                paragraph.font.color.rgb = hex_to_rgb(style_config["body_color"])
        
        # Add accent bar at bottom
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.0),
            Inches(13.333), Inches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(style_config["accent_color"])
        accent_bar.line.fill.background()
    
    # Content slides
    for slide_data in data.get("slides", []):
        layout_idx = 1  # Default: title and content
        
        if slide_data.get("layout") == "section":
            layout_idx = 2  # Section header
        elif slide_data.get("layout") == "two_column":
            layout_idx = 3  # Two content
        elif slide_data.get("layout") == "blank":
            layout_idx = 6  # Blank
        
        try:
            slide_layout = prs.slide_layouts[layout_idx]
        except IndexError:
            slide_layout = prs.slide_layouts[1]
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        if slide.shapes.title and slide_data.get("title"):
            slide.shapes.title.text = slide_data["title"]
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.size = Pt(style_config["title_size"] - 4)
                paragraph.font.name = style_config["title_font"]
                paragraph.font.bold = True
                paragraph.font.color.rgb = hex_to_rgb(style_config["title_color"])
        
        # Bullets
        if slide_data.get("bullets"):
            # Find content placeholder
            content_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_shape = shape
                    break
            
            if content_shape and hasattr(content_shape, "text_frame"):
                tf = content_shape.text_frame
                tf.clear()
                
                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = bullet
                    p.font.size = Pt(style_config["body_size"])
                    p.font.name = style_config["body_font"]
                    p.font.color.rgb = hex_to_rgb(style_config["body_color"])
                    p.level = 0
        
        # Image
        if slide_data.get("image"):
            img_data = slide_data["image"]
            img_path = None
            
            if isinstance(img_data, dict):
                src = img_data.get("src", "")
                if "generate:" in src.lower():
                    prompt = src.split("generate:", 1)[1].strip()
                    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
                    img_path = generate_image(prompt, f"/tmp/pptx_img_{timestamp}.png")
                else:
                    img_path = src
            elif isinstance(img_data, str):
                if "generate:" in img_data.lower():
                    prompt = img_data.split("generate:", 1)[1].strip()
                    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
                    img_path = generate_image(prompt, f"/tmp/pptx_img_{timestamp}.png")
                else:
                    img_path = img_data
            
            if img_path and Path(img_path).exists():
                # Add image to slide
                left = Inches(7)
                top = Inches(1.5)
                width = Inches(5)
                slide.shapes.add_picture(img_path, left, top, width=width)
        
        # Speaker notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]
    
    return prs


def list_templates():
    """List available templates."""
    print("Built-in styles:")
    for style in STYLES:
        print(f"  - {style}")
    
    print("\nCustom templates:")
    if TEMPLATES_DIR.exists():
        templates = list(TEMPLATES_DIR.glob("*.pptx"))
        if templates:
            for t in templates:
                print(f"  - {t.stem}")
        else:
            print("  (none)")
    else:
        print("  (none)")


def save_template(name: str, source_path: str):
    """Save a presentation as a reusable template."""
    TEMPLATES_DIR.mkdir(exist_ok=True)
    
    import shutil
    dest = TEMPLATES_DIR / f"{name}.pptx"
    shutil.copy(source_path, dest)
    print(f"Template saved: {dest}")


def main():
    parser = argparse.ArgumentParser(description="Create PowerPoint presentations")
    
    # Input options
    parser.add_argument("--outline", "-o", help="Markdown outline file")
    parser.add_argument("--json", "-j", help="JSON structure file")
    parser.add_argument("--topic", "-t", help="Topic for AI-generated outline")
    parser.add_argument("--slides", "-n", type=int, default=6, help="Number of slides (for --topic)")
    
    # Template/style options
    parser.add_argument("--template", default="minimal", help="Template name or style")
    parser.add_argument("--list-templates", action="store_true", help="List available templates")
    parser.add_argument("--save-template", metavar="NAME", help="Save presentation as template")
    parser.add_argument("--from", dest="from_file", help="Source file for --save-template")
    
    # Output
    parser.add_argument("--output", "-O", default="presentation.pptx", help="Output file path")
    
    # Features
    parser.add_argument("--generate-images", action="store_true", help="Generate images for slides")
    parser.add_argument("--data-source", help="Data source URI (twenty://, csv, etc.)")
    
    args = parser.parse_args()
    
    # Handle utility commands
    if args.list_templates:
        list_templates()
        return
    
    if args.save_template:
        if not args.from_file:
            print("Error: --save-template requires --from <file>", file=sys.stderr)
            sys.exit(1)
        save_template(args.save_template, args.from_file)
        return
    
    # Determine input
    presentation_data = None
    
    if args.outline:
        presentation_data = parse_outline(args.outline)
    elif args.json:
        with open(args.json, "r") as f:
            presentation_data = json.load(f)
    elif args.topic:
        # Generate outline from topic (placeholder - would use AI)
        presentation_data = {
            "title": args.topic,
            "subtitle": datetime.now().strftime("%B %Y"),
            "slides": [
                {"title": "Overview", "bullets": ["Key point 1", "Key point 2", "Key point 3"]},
                {"title": "Background", "bullets": ["Context", "History", "Current state"]},
                {"title": "Analysis", "bullets": ["Finding 1", "Finding 2", "Finding 3"]},
                {"title": "Recommendations", "bullets": ["Action 1", "Action 2", "Action 3"]},
                {"title": "Next Steps", "bullets": ["Timeline", "Resources", "Follow-up"]},
            ][:args.slides]
        }
        print(f"Note: Generated placeholder outline for '{args.topic}'", file=sys.stderr)
        print("For AI-generated content, provide a detailed outline file.", file=sys.stderr)
    else:
        print("Error: Provide --outline, --json, or --topic", file=sys.stderr)
        sys.exit(1)
    
    # Find template
    template_path = None
    style = "minimal"
    
    if args.template in STYLES:
        style = args.template
    else:
        # Look for custom template
        custom_template = TEMPLATES_DIR / f"{args.template}.pptx"
        if custom_template.exists():
            template_path = str(custom_template)
        elif Path(args.template).exists():
            template_path = args.template
    
    # Create presentation
    prs = create_presentation(presentation_data, style=style, template_path=template_path)
    
    # Save
    output_path = Path(args.output)
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
