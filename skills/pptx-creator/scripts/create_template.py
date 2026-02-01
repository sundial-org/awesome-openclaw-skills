#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "Pillow"]
# ///
"""
Template Creator - Generate branded PowerPoint templates with proper layouts.
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml


from pptx.dml.color import RGBColor

SKILL_DIR = Path(__file__).parent.parent
TEMPLATES_DIR = SKILL_DIR / "templates"

# Template presets with full styling
PRESETS = {
    "minimal": {
        "name": "Minimal",
        "colors": {
            "background": "FFFFFF",
            "title": "1a1a1a",
            "body": "4a4a4a",
            "accent": "0066cc",
            "accent2": "00a86b",
        },
        "fonts": {
            "title": "Helvetica Neue",
            "body": "Helvetica Neue",
        },
        "sizes": {
            "title": 44,
            "subtitle": 24,
            "heading": 36,
            "body": 18,
            "caption": 14,
        }
    },
    "corporate": {
        "name": "Corporate",
        "colors": {
            "background": "FFFFFF",
            "title": "003366",
            "body": "333333",
            "accent": "0066cc",
            "accent2": "ff6600",
        },
        "fonts": {
            "title": "Arial",
            "body": "Arial",
        },
        "sizes": {
            "title": 40,
            "subtitle": 22,
            "heading": 32,
            "body": 18,
            "caption": 12,
        }
    },
    "creative": {
        "name": "Creative",
        "colors": {
            "background": "FAFAFA",
            "title": "2d2d2d",
            "body": "4a4a4a",
            "accent": "ff5722",
            "accent2": "9c27b0",
        },
        "fonts": {
            "title": "Avenir Next",
            "body": "Avenir",
        },
        "sizes": {
            "title": 48,
            "subtitle": 24,
            "heading": 40,
            "body": 20,
            "caption": 14,
        }
    },
    "dark": {
        "name": "Dark Mode",
        "colors": {
            "background": "1a1a2e",
            "title": "ffffff",
            "body": "e0e0e0",
            "accent": "00d4ff",
            "accent2": "ff6b6b",
        },
        "fonts": {
            "title": "SF Pro Display",
            "body": "SF Pro Text",
        },
        "sizes": {
            "title": 44,
            "subtitle": 24,
            "heading": 36,
            "body": 18,
            "caption": 14,
        }
    },
    "executive": {
        "name": "Executive",
        "colors": {
            "background": "FFFFFF",
            "title": "1e3a5f",
            "body": "2c3e50",
            "accent": "c9a227",
            "accent2": "1e3a5f",
        },
        "fonts": {
            "title": "Georgia",
            "body": "Calibri",
        },
        "sizes": {
            "title": 42,
            "subtitle": 22,
            "heading": 34,
            "body": 18,
            "caption": 12,
        }
    },
    "startup": {
        "name": "Startup Pitch",
        "colors": {
            "background": "FFFFFF",
            "title": "2d3436",
            "body": "636e72",
            "accent": "6c5ce7",
            "accent2": "00b894",
        },
        "fonts": {
            "title": "Poppins",
            "body": "Inter",
        },
        "sizes": {
            "title": 48,
            "subtitle": 24,
            "heading": 40,
            "body": 20,
            "caption": 14,
        }
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def create_template(preset_name: str, output_path: Path) -> None:
    """Create a PowerPoint template with proper layouts."""
    preset = PRESETS.get(preset_name, PRESETS["minimal"])
    colors = preset["colors"]
    fonts = preset["fonts"]
    sizes = preset["sizes"]
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # We'll create slides that demonstrate each layout
    # Users can delete these after seeing them
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Presentation Title"
    title.text_frame.paragraphs[0].font.size = Pt(sizes["title"])
    title.text_frame.paragraphs[0].font.name = fonts["title"]
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["title"])
    
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Subtitle goes here\nAuthor Name"
        for para in subtitle.text_frame.paragraphs:
            para.font.size = Pt(sizes["subtitle"])
            para.font.name = fonts["body"]
            para.font.color.rgb = hex_to_rgb(colors["body"])
    
    # Add accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8),
        Inches(13.333), Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(colors["accent"])
    shape.line.fill.background()
    
    # 2. Section Header
    slide_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = "Section Title"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.size = Pt(sizes["title"])
            para.font.name = fonts["title"]
            para.font.bold = True
            para.font.color.rgb = hex_to_rgb(colors["accent"])
    
    # 3. Title and Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = "Slide Title"
    for para in slide.shapes.title.text_frame.paragraphs:
        para.font.size = Pt(sizes["heading"])
        para.font.name = fonts["title"]
        para.font.color.rgb = hex_to_rgb(colors["title"])
    
    # Find content placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            
            bullets = ["First key point", "Second key point", "Third key point with more detail", "Final point"]
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(sizes["body"])
                p.font.name = fonts["body"]
                p.font.color.rgb = hex_to_rgb(colors["body"])
                p.level = 0
            break
    
    # 4. Two Column Layout
    slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = "Two Column Layout"
        for para in slide.shapes.title.text_frame.paragraphs:
            para.font.size = Pt(sizes["heading"])
            para.font.name = fonts["title"]
            para.font.color.rgb = hex_to_rgb(colors["title"])
    
    # 5. Image with Caption (simulated)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Image Slide"
    for para in slide.shapes.title.text_frame.paragraphs:
        para.font.size = Pt(sizes["heading"])
        para.font.name = fonts["title"]
        para.font.color.rgb = hex_to_rgb(colors["title"])
    
    # Add placeholder rectangle for image
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.8),
        Inches(7), Inches(5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = hex_to_rgb(colors["body"])
    img_placeholder.fill.fore_color.brightness = 0.8
    img_placeholder.line.color.rgb = hex_to_rgb(colors["accent"])
    
    # Add text box for caption
    caption_box = slide.shapes.add_textbox(
        Inches(8.5), Inches(2),
        Inches(4), Inches(4)
    )
    tf = caption_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Image description and supporting text goes here. This area can contain bullet points or paragraphs."
    p.font.size = Pt(sizes["body"])
    p.font.name = fonts["body"]
    p.font.color.rgb = hex_to_rgb(colors["body"])
    
    # 6. Quote/Callout Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1])
    
    # Add large quote
    quote_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2),
        Inches(10), Inches(3)
    )
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"This is a powerful quote or key insight that you want to emphasize."'
    p.font.size = Pt(36)
    p.font.name = fonts["title"]
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(colors["accent"])
    p.alignment = PP_ALIGN.CENTER
    
    # Attribution
    p2 = tf.add_paragraph()
    p2.text = "— Attribution"
    p2.font.size = Pt(sizes["body"])
    p2.font.name = fonts["body"]
    p2.font.color.rgb = hex_to_rgb(colors["body"])
    p2.alignment = PP_ALIGN.CENTER
    
    # 7. Thank You / End Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title = slide.shapes.title
    title.text = "Thank You"
    for para in title.text_frame.paragraphs:
        para.font.size = Pt(sizes["title"])
        para.font.name = fonts["title"]
        para.font.bold = True
        para.font.color.rgb = hex_to_rgb(colors["accent"])
        para.alignment = PP_ALIGN.CENTER
    
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Questions?\ncontact@example.com"
        for para in subtitle.text_frame.paragraphs:
            para.font.size = Pt(sizes["subtitle"])
            para.font.name = fonts["body"]
            para.font.color.rgb = hex_to_rgb(colors["body"])
            para.alignment = PP_ALIGN.CENTER
    
    # Add accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(colors["accent"])
    shape.line.fill.background()
    
    # Save
    prs.save(output_path)
    print(f"Created template: {output_path}")


def create_all_templates():
    """Generate all preset templates."""
    TEMPLATES_DIR.mkdir(exist_ok=True)
    
    for preset_name in PRESETS:
        output_path = TEMPLATES_DIR / f"{preset_name}.pptx"
        create_template(preset_name, output_path)
    
    # Also create a config file with template metadata
    config = {
        "templates": {
            name: {
                "name": preset["name"],
                "description": f"{preset['name']} theme with {preset['fonts']['title']} font",
                "colors": preset["colors"],
            }
            for name, preset in PRESETS.items()
        }
    }
    
    config_path = TEMPLATES_DIR / "templates.json"
    with open(config_path, "w") as f:
        json.dump(config, f, indent=2)
    print(f"Created config: {config_path}")


def list_presets():
    """List available template presets."""
    print("Available template presets:\n")
    for name, preset in PRESETS.items():
        colors = preset["colors"]
        print(f"  {name}")
        print(f"    Name: {preset['name']}")
        print(f"    Fonts: {preset['fonts']['title']} / {preset['fonts']['body']}")
        print(f"    Accent: #{colors['accent']}")
        print()


def main():
    parser = argparse.ArgumentParser(description="Create PowerPoint templates")
    parser.add_argument("--preset", "-p", help="Preset name to generate")
    parser.add_argument("--all", "-a", action="store_true", help="Generate all templates")
    parser.add_argument("--list", "-l", action="store_true", help="List available presets")
    parser.add_argument("--output", "-o", help="Output path (default: templates/<preset>.pptx)")
    
    args = parser.parse_args()
    
    if args.list:
        list_presets()
        return
    
    if args.all:
        create_all_templates()
        return
    
    if args.preset:
        if args.preset not in PRESETS:
            print(f"Unknown preset: {args.preset}")
            print(f"Available: {', '.join(PRESETS.keys())}")
            return
        
        TEMPLATES_DIR.mkdir(exist_ok=True)
        output_path = Path(args.output) if args.output else TEMPLATES_DIR / f"{args.preset}.pptx"
        create_template(args.preset, output_path)
        return
    
    parser.print_help()


if __name__ == "__main__":
    main()
